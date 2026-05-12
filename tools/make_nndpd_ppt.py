# -*- coding: utf-8 -*-
"""生成 NN-DPD 离线训练示例的汇报 PPT。

说明：基于 MathWorks 官方示例
"Neural Network for Digital Predistortion Design - Offline Training"
把这张频谱对比图的产出流程做成一份向老师汇报的 PPT。
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 全局样式 ----------
TITLE_FONT = "微软雅黑"
BODY_FONT = "微软雅黑"
MONO_FONT = "Consolas"

COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # 深蓝
COLOR_ACCENT = RGBColor(0xE8, 0x7A, 0x00)       # 橙
COLOR_YELLOW = RGBColor(0xC9, 0xA6, 0x27)
COLOR_GRAY = RGBColor(0x55, 0x55, 0x55)
COLOR_LIGHT = RGBColor(0xE8, 0xEE, 0xF4)


def set_cn_font(run, name=BODY_FONT):
    """给 run 设定中英文字体。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    # 东亚字符字体
    for tag in ("eastAsia", "ascii", "hAnsi"):
        e = rPr.find(qn(f"a:{tag}"))
        if e is None:
            from lxml import etree
            e = etree.SubElement(rPr, qn(f"a:{tag}"))
        e.set("typeface", name)


def add_title(slide, text, color=COLOR_PRIMARY, size=30, left=Inches(0.5), top=Inches(0.35)):
    tx = slide.shapes.add_textbox(left, top, Inches(12.33), Inches(0.9))
    tf = tx.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    set_cn_font(r, TITLE_FONT)
    # 加一条底色横线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.33), Emu(30000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()
    return tx


def add_footer(slide, text="NN-DPD 离线训练 · 汇报材料"):
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = COLOR_GRAY
    set_cn_font(r, BODY_FONT)


def add_bullets(slide, items, left=Inches(0.7), top=Inches(1.5), width=Inches(12.0),
                height=Inches(5.3), size=18, line_spacing=1.25):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(6)

    for i, item in enumerate(items):
        # item 可以是 str，或者 (level, text)
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        bullet = "• " if level == 0 else "– "
        r = p.add_run()
        r.text = bullet + text
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = COLOR_GRAY if level > 0 else RGBColor(0x22, 0x22, 0x22)
        set_cn_font(r, BODY_FONT)
    return tx


def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    # 把背景挪到最底层
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def add_colored_box(slide, left, top, width, height, fill, text, size=16, bold=False,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    set_cn_font(r, BODY_FONT)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=COLOR_PRIMARY, weight=2.25):
    connector = slide.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    # 加箭头
    lnEl = connector.line._get_or_add_ln()
    from lxml import etree
    tail = etree.SubElement(lnEl, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return connector


# ---------- 构造 PPT ----------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1 封面 ---
    s = add_blank_slide(prs)
    # 大色块
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), prs.slide_width, Inches(2.6))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    tx = s.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(12), Inches(1.2))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "神经网络数字预失真(NN-DPD) 离线训练示例"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    set_cn_font(r, TITLE_FONT)

    sub = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(12), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "——那张 PA 输出频谱对比图是如何得出的？"
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
    set_cn_font(r, TITLE_FONT)

    meta = s.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(12), Inches(1.5))
    tf = meta.text_frame
    for line in [
        "来源：MathWorks 官方示例 Neural Network for Digital Predistortion Design - Offline Training",
        "文献基础：[1] Indirect Learning, ISWCS 2008  ·  [2] ARVTDNN, IEEE TNNLS 2019",
        "汇报用途：说明图的数据来源、网络结构、训练流程与指标对比",
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = COLOR_GRAY
        set_cn_font(r, BODY_FONT)
    add_footer(s, "ZTE 项目 · DPD 专题")

    # --- Slide 2 图是什么 ---
    s = add_blank_slide(prs)
    add_title(s, "这张图展示了什么")
    add_bullets(s, [
        "横轴：频率（基带归一化，单位 GHz 等价刻度；实际采样率 614.4 MHz）",
        "纵轴：功率谱密度(dBm)，spectrumAnalyzer 做 16 次谱平均得到平滑曲线",
        "三条曲线 = 同一段 100 MHz OFDM 信号在三种 DPD 方案下通过 PA 后的输出频谱",
        (1, "蓝色 No DPD：信号直接送 PA，暴露 PA 非线性"),
        (1, "橙色 Memory Polynomial DPD：传统交叉项记忆多项式预失真"),
        (1, "黄色 Neural Network DPD：训练好的神经网络预失真"),
        "看点：主瓣两侧的「带外频谱再生 (spectral regrowth)」越低越好",
        (1, "蓝 → 橙 → 黄 肩部逐级压低，NN-DPD 抑制带外泄漏最明显"),
    ])
    add_footer(s)

    # --- Slide 3 整体工作流 ---
    s = add_blank_slide(prs)
    add_title(s, "整体工作流：间接学习法 (Indirect Learning)")

    # 训练路径
    tbl_y = Inches(1.7)
    add_colored_box(s, Inches(0.6), tbl_y, Inches(1.9), Inches(0.8), COLOR_PRIMARY, "训练信号 u(n)")
    add_colored_box(s, Inches(2.9), tbl_y, Inches(1.6), Inches(0.8), COLOR_ACCENT, "PA 采样")
    add_colored_box(s, Inches(4.9), tbl_y, Inches(1.9), Inches(0.8), COLOR_PRIMARY, "PA 输出 x(n)")
    add_colored_box(s, Inches(7.2), tbl_y, Inches(2.3), Inches(0.8), COLOR_YELLOW, "NN 输入=x, 目标=u", color=RGBColor(0,0,0))
    add_colored_box(s, Inches(9.9), tbl_y, Inches(3.0), Inches(0.8), COLOR_PRIMARY, "得到 PA 的逆函数")

    add_arrow(s, Inches(2.5), Inches(2.1), Inches(2.9), Inches(2.1))
    add_arrow(s, Inches(4.5), Inches(2.1), Inches(4.9), Inches(2.1))
    add_arrow(s, Inches(6.8), Inches(2.1), Inches(7.2), Inches(2.1))
    add_arrow(s, Inches(9.5), Inches(2.1), Inches(9.9), Inches(2.1))

    lab = s.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(6), Inches(0.3))
    r = lab.text_frame.paragraphs[0].add_run()
    r.text = "① 训练阶段（离线，采集 (u, x) 数据对）"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = COLOR_PRIMARY
    set_cn_font(r)

    # 部署路径
    tbl_y = Inches(3.6)
    add_colored_box(s, Inches(0.6), tbl_y, Inches(1.9), Inches(0.8), COLOR_PRIMARY, "发射信号 u(n)")
    add_colored_box(s, Inches(2.9), tbl_y, Inches(2.3), Inches(0.8), COLOR_YELLOW, "训练好的 NN-DPD", color=RGBColor(0,0,0))
    add_colored_box(s, Inches(5.6), tbl_y, Inches(1.9), Inches(0.8), COLOR_PRIMARY, "预失真信号 y(n)")
    add_colored_box(s, Inches(7.9), tbl_y, Inches(1.6), Inches(0.8), COLOR_ACCENT, "PA")
    add_colored_box(s, Inches(9.9), tbl_y, Inches(3.0), Inches(0.8), COLOR_PRIMARY, "线性化输出 z(n)")

    add_arrow(s, Inches(2.5), Inches(4.0), Inches(2.9), Inches(4.0))
    add_arrow(s, Inches(5.2), Inches(4.0), Inches(5.6), Inches(4.0))
    add_arrow(s, Inches(7.5), Inches(4.0), Inches(7.9), Inches(4.0))
    add_arrow(s, Inches(9.5), Inches(4.0), Inches(9.9), Inches(4.0))

    lab = s.shapes.add_textbox(Inches(0.6), Inches(3.25), Inches(6), Inches(0.3))
    r = lab.text_frame.paragraphs[0].add_run()
    r.text = "② 部署阶段（把 NN-DPD 串在 PA 前面）"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = COLOR_PRIMARY
    set_cn_font(r)

    add_bullets(s, [
        "核心思想：NN 学到的是 PA 的逆；放到 PA 前面，预失真 × PA = 线性",
        "优点：避免在线拟合 PA 模型，训练完直接推理；对硬件友好",
    ], top=Inches(4.9), size=16)

    add_footer(s)

    # --- Slide 4 数据生成 ---
    s = add_blank_slide(prs)
    add_title(s, "步骤 1：数据生成（100 MHz OFDM 激励）")
    add_bullets(s, [
        "基于 5G 参数生成 OFDM 信号（helperOFDMParameters）",
        (1, "子载波间隔 SCS = 30 kHz，FFT 长度 = 4096"),
        (1, "数据子载波 3276 个，CP 长度 288，每帧 6 个 OFDM 符号"),
        (1, "子载波上调制 16-QAM 符号"),
        "过采样因子 osf = 5 → 采样率 30 kHz × 4096 × 5 ≈ 614.4 MHz",
        (1, "目的：抓住 PA 产生的高阶非线性分量（带外频谱）"),
        "生成三套相互独立的数据（相同分布，不同随机序列）",
        (1, "txWaveTrain / txWaveVal / txWaveTest"),
        "全部送入 helperNNDPDPowerAmplifier 获得 PA 输出，得到 (u, x) 对",
        (1, "数据源可选：真实 NI VST + NXP Doherty PA / 仿真 PA 模型 / 已保存数据"),
    ])
    add_footer(s)

    # --- Slide 5 特征构造：ARVTDNN ---
    s = add_blank_slide(prs)
    add_title(s, "步骤 2：特征构造（ARVTDNN 的「专家知识」输入）")
    add_bullets(s, [
        "背景：PA 记忆多项式模型",
        (1, "x(n) = ΣΣ c_{m,k} · u(n−m) · |u(n−m)|^k   (m=0..M−1, k=0..K−1)"),
        "把记忆多项式的先验打包成神经网络输入，降低网络规模",
        "取 M=5（记忆深度），K=5（非线性阶数）",
        "每个时刻 n 的输入特征：",
        (1, "I/Q 分量及其 M 个延迟：2M = 10 维"),
        (1, "幅度高阶项 |u(n−m)|^k，k=2..5：(K−1)·M = 20 维"),
        (1, "合计输入维度 = 30"),
        "目标输出（训练标签）：对应 PA 输入 u(n) 的 I/Q → 2 维",
    ], size=18)
    add_footer(s)

    # --- Slide 6 网络结构 ---
    s = add_blank_slide(prs)
    add_title(s, "步骤 3：网络结构（全连接 TDNN）")

    y = Inches(2.0)
    layers = [
        ("输入层\n30 维", COLOR_PRIMARY),
        ("FC 30\n+LeakyReLU", COLOR_ACCENT),
        ("FC 24\n+LeakyReLU", COLOR_ACCENT),
        ("FC 19\n+LeakyReLU", COLOR_ACCENT),
        ("输出层\n2 维 (I,Q)", COLOR_YELLOW),
    ]
    x = Inches(0.7)
    w = Inches(2.1)
    gap = Inches(0.4)
    for i, (name, c) in enumerate(layers):
        color_txt = RGBColor(0, 0, 0) if c == COLOR_YELLOW else RGBColor(0xFF, 0xFF, 0xFF)
        add_colored_box(s, x, y, w, Inches(1.2), c, name, size=14, bold=True, color=color_txt)
        if i < len(layers) - 1:
            add_arrow(s, x + w, y + Inches(0.6), x + w + gap, y + Inches(0.6))
        x = x + w + gap

    add_bullets(s, [
        "神经元数量按 0.8 的比例逐层缩减（30 → 24 → 19）",
        "激活函数：LeakyReLU(0.01)；输出层为线性",
        "输出代表预失真后信号的实部/虚部",
        "网络规模小，适合实时推理（推理每样本仅 ~2k MAC）",
    ], top=Inches(3.8), size=18)

    add_footer(s)

    # --- Slide 7 训练配置 ---
    s = add_blank_slide(prs)
    add_title(s, "步骤 4：训练配置（离线监督学习）")
    add_bullets(s, [
        "优化器：Adam；损失函数：MSE（回归 I/Q 两路输出）",
        "Mini-batch = 1024，最大 epoch = 200",
        "初始学习率 4e-4；每 5 个 epoch 乘 0.95（piecewise）",
        "验证频率：每 epoch 两次；ValidationPatience = 5 → 早停",
        "输出选择：OutputNetwork = best-validation-loss",
        "训练硬件：Intel Xeon W-2133 CPU 约 6 分钟收敛",
        "示例默认加载预训练权重：nndpdIn30Fact08.mat",
        (1, "更换 PA / 带宽 / 目标功率时需要重新训练"),
    ])
    add_footer(s)

    # --- Slide 8 测试与绘图 ---
    s = add_blank_slide(prs)
    add_title(s, "步骤 5：测试流程（图就是在这里生成的）")
    add_bullets(s, [
        "用测试信号 txWaveTest 走三条支路进 PA，得到三组 PA 输出：",
        (1, "paOutputTest：直接送 PA（No DPD）"),
        (1, "paOutputMP：经记忆多项式 DPD 预失真后再送 PA"),
        (1, "paOutputNN：经 NN-DPD 预失真后再送 PA"),
        "关键代码：",
    ], top=Inches(1.5))

    code_tx = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12.0), Inches(2.6))
    tf = code_tx.text_frame
    tf.word_wrap = True
    code_lines = [
        "dpdOutNN  = predict(netDPD, inputMtxTest);",
        "paOutputNN = pa( complex(dpdOutNN(:,1), dpdOutNN(:,2)) / scalingFactor );",
        "dpdOutMP  = helperNNDPDMemoryPolynomial(txWaveTest, txWaveTrain, paOutputTrain, 5, 5);",
        "paOutputMP = pa(dpdOutMP);",
        "sa = helperPACharPlotSpectrum([paOutputTest paOutputMP paOutputNN], ...",
        "       {'No DPD','Memory Polynomial DPD','Neural Network DPD'}, ...",
        "       ofdmParams.OversamplingFactor, \"Modulated\", [-130 -50]);",
    ]
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.name = MONO_FONT
        r.font.color.rgb = RGBColor(0x11, 0x35, 0x5D)

    add_footer(s)

    # --- Slide 9 结果数值 ---
    s = add_blank_slide(prs)
    add_title(s, "结果对比：频谱图与三项定量指标互相印证")

    # 画表格
    rows = 4
    cols = 4
    left = Inches(0.8)
    top = Inches(1.7)
    width = Inches(11.7)
    height = Inches(2.2)
    table = s.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["方案", "ACPR (dB) ↓", "NMSE (dB) ↓", "EVM (%) ↓"]
    data = [
        ["No DPD", "-28.674", "-21.287", "6.868"],
        ["Memory Polynomial DPD", "-33.889", "-27.984", "2.823"],
        ["Neural Network DPD", "-38.886", "-33.423", "1.568"],
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        set_cn_font(r)
        cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_PRIMARY

    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(15)
            # NN-DPD 行高亮
            if i == 3:
                r.font.bold = True
                r.font.color.rgb = COLOR_ACCENT
            else:
                r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            set_cn_font(r)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT if i % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)

    add_bullets(s, [
        "三个指标都是越小越好；NN-DPD 相对 No DPD，ACPR 改善约 10 dB，EVM 从 6.87% 降到 1.57%",
        "频谱图的肩部高度与 ACPR 数值一一对应，图与表相互佐证",
        "结论：对于 100 MHz OFDM 信号 + NXP Doherty PA，NN-DPD 优于传统记忆多项式 DPD",
    ], top=Inches(4.2), size=17)
    add_footer(s)

    # --- Slide 10 汇报要点总结 ---
    s = add_blank_slide(prs)
    add_title(s, "汇报要点 · 一句话版本")
    tx = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("这张图展示了同一段 100 MHz OFDM 信号经过 PA 后，在三种情况（无 DPD / 记忆多项式 DPD / NN-DPD）"
              "下的功率谱，用于直观对比 DPD 对 PA 非线性导致的带外频谱再生的抑制能力；结果显示 NN-DPD 的带外"
              "泄漏最低，印证了 ACPR 由 −28.7 dB 改善到 −38.9 dB 的定量结论。")
    r.font.size = Pt(20); r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    set_cn_font(r)

    add_title(s, "展开讲的顺序（建议）", top=Inches(3.8), size=22)
    add_bullets(s, [
        "① 为什么做 DPD：PA 非线性 → 带外频谱再生 → 不满足发射机规范",
        "② 间接学习法：用 (u, x) 训练 NN 学 PA 的逆",
        "③ ARVTDNN 输入：I/Q 延迟 + 幅度多阶项（记忆多项式先验）",
        "④ 训练：Adam + MSE，离线监督学习，早停",
        "⑤ 验证：同一测试信号走三条支路 → 测 ACPR/NMSE/EVM + 画频谱",
        "⑥ 结论：图（定性） + 表（定量）共同说明 NN-DPD 的优势",
    ], top=Inches(4.4), size=17)

    add_footer(s, "汇报完毕 · Q&A")

    out = Path(__file__).resolve().parent.parent / "NN_DPD_离线训练_汇报.pptx"
    prs.save(out)
    print(f"PPT 已生成: {out}")


if __name__ == "__main__":
    build()
